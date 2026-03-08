#!/usr/bin/env python3
"""
Generate a table stress-test PPTX that exercises table rendering thoroughly.

Creates slides covering simple tables, merged cells, text alignments,
banded rows, and nested text formatting.

Usage:
    python3 scripts/generate-table-stress-test.py
    # Output: test-data/table-stress-test.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ROOT = __import__("pathlib").Path(__file__).resolve().parent.parent

SLIDE_WIDTH = Emu(12192000)   # 13.333 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches


def add_label(slide, x, y, w, h, text, font_size=10):
    """Add a text label to the slide."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def set_cell_fill(cell, color):
    """Set a solid fill on a table cell. color is an RGBColor."""
    tcPr = cell._tc.get_or_add_tcPr()
    # Remove existing fill
    for child in list(tcPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "gradFill", "noFill"):
            tcPr.remove(child)
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {
        "val": str(color)
    })
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def set_cell_borders(cell, color=RGBColor(0x00, 0x00, 0x00), width_pt=1):
    """Set borders on a table cell."""
    tcPr = cell._tc.get_or_add_tcPr()
    width_emu = int(width_pt * 12700)
    color_str = str(color)

    for border_name in ["lnL", "lnR", "lnT", "lnB"]:
        # Remove existing
        for child in list(tcPr):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == border_name:
                tcPr.remove(child)

        ln = tcPr.makeelement(qn(f"a:{border_name}"), {"w": str(width_emu)})
        solidFill = ln.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": color_str})
        solidFill.append(srgbClr)
        ln.append(solidFill)
        tcPr.append(ln)


def slide1_simple_table(prs):
    """Slide 1: Simple 4x4 table with borders and cell fills."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 1: Simple 4x4 Table with Borders and Cell Fills", 18)

    rows, cols = 4, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(1.0), Inches(1.0), Inches(10), Inches(5.0))
    table = table_shape.table

    # Header row
    header_color = RGBColor(0x00, 0x52, 0x8A)
    headers = ["Product", "Q1 Sales", "Q2 Sales", "Q3 Sales"]
    for c in range(cols):
        cell = table.cell(0, c)
        cell.text = headers[c]
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(14)
                run.font.name = "Calibri"
        set_cell_fill(cell, header_color)
        set_cell_borders(cell, RGBColor(0x00, 0x33, 0x55), 2)

    # Data rows
    data = [
        ["Widget A", "$1,250", "$1,430", "$1,890"],
        ["Widget B", "$890", "$1,050", "$1,200"],
        ["Widget C", "$2,100", "$1,950", "$2,350"],
    ]
    row_colors = [RGBColor(0xE8, 0xF0, 0xFE), RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xE8, 0xF0, 0xFE)]

    for r in range(3):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            cell.text = data[r][c]
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.name = "Calibri"
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            set_cell_fill(cell, row_colors[r])
            set_cell_borders(cell, RGBColor(0x99, 0x99, 0x99), 1)


def slide2_merged_cells(prs):
    """Slide 2: Merged cells (horizontal and vertical spans)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 2: Merged Cells (Horizontal and Vertical Spans)", 18)

    rows, cols = 5, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(1.0), Inches(1.0), Inches(10), Inches(5.5))
    table = table_shape.table

    # Horizontal merge: row 0, cols 0-3 (full width header)
    cell_00 = table.cell(0, 0)
    cell_03 = table.cell(0, 3)
    cell_00.merge(cell_03)
    cell_00.text = "Quarterly Revenue Report (Horizontally Merged)"
    for p in cell_00.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    set_cell_fill(cell_00, RGBColor(0x1A, 0x1A, 0x6B))

    # Vertical merge: col 0, rows 1-2
    cell_10 = table.cell(1, 0)
    cell_20 = table.cell(2, 0)
    cell_10.merge(cell_20)
    cell_10.text = "North Region\n(Vertically Merged)"
    for p in cell_10.text_frame.paragraphs:
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.name = "Calibri"
    set_cell_fill(cell_10, RGBColor(0xD4, 0xE6, 0xF1))

    # Vertical merge: col 0, rows 3-4
    cell_30 = table.cell(3, 0)
    cell_40 = table.cell(4, 0)
    cell_30.merge(cell_40)
    cell_30.text = "South Region\n(Vertically Merged)"
    for p in cell_30.text_frame.paragraphs:
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.name = "Calibri"
    set_cell_fill(cell_30, RGBColor(0xFA, 0xDB, 0xD8))

    # Fill in regular data cells
    regular_data = {
        (1, 1): "Q1: $500K", (1, 2): "Q2: $620K", (1, 3): "Q3: $710K",
        (2, 1): "Q1: $380K", (2, 2): "Q2: $410K", (2, 3): "Q3: $450K",
        (3, 1): "Q1: $290K", (3, 2): "Q2: $310K", (3, 3): "Q3: $365K",
        (4, 1): "Q1: $195K", (4, 2): "Q2: $220K", (4, 3): "Q3: $240K",
    }
    for (r, c), text in regular_data.items():
        cell = table.cell(r, c)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.name = "Calibri"
        set_cell_borders(cell, RGBColor(0x99, 0x99, 0x99), 1)


def slide3_text_alignments(prs):
    """Slide 3: Tables with different cell text alignments."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 3: Cell Text Alignments (3x3 Grid: V x H)", 18)

    rows, cols = 4, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.0), Inches(9), Inches(5.5))
    table = table_shape.table

    # Set column widths
    for c in range(cols):
        table.columns[c].width = Inches(2.25)

    # Row heights
    for r in range(rows):
        table.rows[r].height = Inches(1.375)

    # Header row labels
    h_labels = ["", "Left", "Center", "Right"]
    v_labels = ["", "Top", "Middle", "Bottom"]

    h_aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT]
    v_anchors = [MSO_ANCHOR.TOP, MSO_ANCHOR.TOP, MSO_ANCHOR.MIDDLE, MSO_ANCHOR.BOTTOM]

    for c in range(cols):
        cell = table.cell(0, c)
        cell.text = h_labels[c]
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.name = "Calibri"
        set_cell_fill(cell, RGBColor(0xD5, 0xE8, 0xD4))

    for r in range(1, rows):
        cell = table.cell(r, 0)
        cell.text = v_labels[r]
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.name = "Calibri"
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_cell_fill(cell, RGBColor(0xD5, 0xE8, 0xD4))

    # Fill the 3x3 alignment grid
    for r in range(1, 4):
        for c in range(1, 4):
            cell = table.cell(r, c)
            cell.text = f"{v_labels[r]}-{h_labels[c]}"
            cell.vertical_anchor = v_anchors[r]
            for p in cell.text_frame.paragraphs:
                p.alignment = h_aligns[c]
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.name = "Calibri"
            set_cell_borders(cell, RGBColor(0x66, 0x66, 0x66), 1)


def slide4_banded_rows(prs):
    """Slide 4: Table with alternating row colors (banded rows)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 4: Alternating Row Colors (Banded Rows)", 18)

    rows, cols = 8, 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.0), Inches(11.5), Inches(5.5))
    table = table_shape.table

    # Header
    headers = ["ID", "Name", "Department", "Start Date", "Salary"]
    header_color = RGBColor(0x2C, 0x3E, 0x50)
    for c in range(cols):
        cell = table.cell(0, c)
        cell.text = headers[c]
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.name = "Calibri"
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        set_cell_fill(cell, header_color)

    # Data with alternating colors
    data = [
        ["001", "Alice Johnson", "Engineering", "2021-03-15", "$95,000"],
        ["002", "Bob Smith", "Marketing", "2020-07-22", "$72,000"],
        ["003", "Carol Williams", "Engineering", "2022-01-10", "$88,000"],
        ["004", "David Brown", "Finance", "2019-11-05", "$110,000"],
        ["005", "Eve Davis", "HR", "2023-06-18", "$65,000"],
        ["006", "Frank Miller", "Engineering", "2021-09-01", "$92,000"],
        ["007", "Grace Wilson", "Marketing", "2022-04-12", "$78,000"],
    ]

    band_colors = [RGBColor(0xEC, 0xF0, 0xF1), RGBColor(0xFF, 0xFF, 0xFF)]

    for r, row_data in enumerate(data):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            cell.text = row_data[c]
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.name = "Calibri"
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            set_cell_fill(cell, band_colors[r % 2])
            set_cell_borders(cell, RGBColor(0xBD, 0xC3, 0xC7), 1)


def slide5_nested_formatting(prs):
    """Slide 5: Nested text formatting in cells (bold, italic, bullets)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 5: Rich Text Formatting in Table Cells", 18)

    rows, cols = 3, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.0), Inches(11.5), Inches(5.5))
    table = table_shape.table

    # Set wider row heights
    for r in range(rows):
        table.rows[r].height = Inches(1.8)

    # Cell (0,0): Bold and italic mixed
    cell = table.cell(0, 0)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Bold text "
    r1.font.bold = True
    r1.font.size = Pt(12)
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = "and italic text "
    r2.font.italic = True
    r2.font.size = Pt(12)
    r2.font.name = "Calibri"
    r3 = p.add_run()
    r3.text = "and bold italic"
    r3.font.bold = True
    r3.font.italic = True
    r3.font.size = Pt(12)
    r3.font.name = "Calibri"

    # Cell (0,1): Different font sizes
    cell = table.cell(0, 1)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    sizes = [8, 10, 12, 16, 20]
    for s in sizes:
        r = p.add_run()
        r.text = f"{s}pt "
        r.font.size = Pt(s)
        r.font.name = "Calibri"

    # Cell (0,2): Colored text
    cell = table.cell(0, 2)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    colors = [
        ("Red ", RGBColor(0xFF, 0x00, 0x00)),
        ("Green ", RGBColor(0x00, 0x80, 0x00)),
        ("Blue ", RGBColor(0x00, 0x00, 0xFF)),
        ("Orange ", RGBColor(0xFF, 0xA5, 0x00)),
        ("Purple", RGBColor(0x80, 0x00, 0x80)),
    ]
    for text, color in colors:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(14)
        r.font.name = "Calibri"
        r.font.color.rgb = color
        r.font.bold = True

    # Cell (1,0): Multiple paragraphs
    cell = table.cell(1, 0)
    tf = cell.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r = p1.add_run()
    r.text = "First paragraph"
    r.font.size = Pt(12)
    r.font.name = "Calibri"
    r.font.bold = True

    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = "Second paragraph with normal text"
    r.font.size = Pt(11)
    r.font.name = "Calibri"

    p3 = tf.add_paragraph()
    r = p3.add_run()
    r.text = "Third paragraph in italics"
    r.font.size = Pt(11)
    r.font.name = "Calibri"
    r.font.italic = True

    # Cell (1,1): Underline styles
    cell = table.cell(1, 1)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Underlined text"
    r.font.size = Pt(12)
    r.font.name = "Calibri"
    r.font.underline = True

    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = "Bold underlined"
    r.font.size = Pt(12)
    r.font.name = "Calibri"
    r.font.bold = True
    r.font.underline = True

    p3 = tf.add_paragraph()
    r = p3.add_run()
    r.text = "Normal after underline"
    r.font.size = Pt(12)
    r.font.name = "Calibri"

    # Cell (1,2): Mixed fonts
    cell = table.cell(1, 2)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    fonts = [("Calibri ", "Calibri"), ("Arial ", "Arial"), ("Courier ", "Courier New")]
    for text, font in fonts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(12)
        r.font.name = font

    # Cell (2,0): Bullet-like content (using dash prefix)
    cell = table.cell(2, 0)
    tf = cell.text_frame
    tf.word_wrap = True
    items = ["First item", "Second item", "Third item", "Fourth item"]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"\u2022 {item}"
        r.font.size = Pt(11)
        r.font.name = "Calibri"

    # Cell (2,1): Numbers list
    cell = table.cell(2, 1)
    tf = cell.text_frame
    tf.word_wrap = True
    for i in range(4):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"{i+1}. Step number {i+1}"
        r.font.size = Pt(11)
        r.font.name = "Calibri"

    # Cell (2,2): Long wrapped text
    cell = table.cell(2, 2)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = (
        "This cell contains a longer passage of text that should wrap within "
        "the cell boundaries. It tests how the renderer handles text overflow "
        "and word wrapping inside table cells."
    )
    r.font.size = Pt(11)
    r.font.name = "Calibri"

    # Set borders on all cells
    for r in range(rows):
        for c in range(cols):
            set_cell_borders(table.cell(r, c), RGBColor(0x66, 0x66, 0x66), 1)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide1_simple_table(prs)
    slide2_merged_cells(prs)
    slide3_text_alignments(prs)
    slide4_banded_rows(prs)
    slide5_nested_formatting(prs)

    output_path = ROOT / "test-data" / "table-stress-test.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"Generated: {output_path}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
