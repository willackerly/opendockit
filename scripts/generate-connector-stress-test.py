#!/usr/bin/env python3
"""
Generate a connector stress-test PPTX that exercises connector rendering.

Creates slides covering straight connectors with arrow styles, bent (elbow)
connectors, curved connectors, and different line styles.

Usage:
    python3 scripts/generate-connector-stress-test.py
    # Output: test-data/connector-stress-test.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

ROOT = __import__("pathlib").Path(__file__).resolve().parent.parent

SLIDE_WIDTH = Emu(12192000)   # 13.333 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches


def add_label(slide, x, y, w, h, text, font_size=10, color=None):
    """Add a text label to the slide."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = "Calibri"
    run.font.color.rgb = color or RGBColor(0x33, 0x33, 0x33)
    return txBox


def add_endpoint_dot(slide, x_inches, y_inches, color=(0x99, 0x99, 0x99)):
    """Add a small circle to mark a connector endpoint."""
    dot = slide.shapes.add_shape(
        9,  # OVAL
        Inches(x_inches) - Emu(36000),
        Inches(y_inches) - Emu(36000),
        Emu(72000), Emu(72000)
    )
    spPr = dot._element.spPr
    for child in list(spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "gradFill", "noFill"):
            spPr.remove(child)
    solidFill = spPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {
        "val": "%02X%02X%02X" % (color[0], color[1], color[2])
    })
    solidFill.append(srgbClr)
    spPr.append(solidFill)
    # Remove outline
    for child in list(spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "ln":
            spPr.remove(child)
    return dot


def add_connector(slide, x1, y1, x2, y2, connector_type="straight",
                  color=(0x00, 0x00, 0x00), width_pt=1.5,
                  head_end=None, tail_end=None,
                  dash_style=None):
    """
    Add a connector shape between two points.
    connector_type: "straight", "bent", or "curved"
    head_end/tail_end: arrow type string or None
    dash_style: "solid", "dash", "dot", "dashDot", "lgDash", "lgDashDot" etc.
    """
    # Connector type mapping for cxnSp
    type_map = {
        "straight": "line",
        "bent": "bentConnector3",
        "curved": "curvedConnector3",
    }
    prst = type_map.get(connector_type, "line")

    # Calculate position and size
    left = min(Inches(x1), Inches(x2))
    top = min(Inches(y1), Inches(y2))
    width = abs(Inches(x2) - Inches(x1))
    height = abs(Inches(y2) - Inches(y1))

    # Determine if we need to flip
    flipH = "1" if x2 < x1 else "0"
    flipV = "1" if y2 < y1 else "0"

    # Ensure minimum size
    if width == 0:
        width = Emu(1)
    if height == 0:
        height = Emu(1)

    # Build cxnSp element
    nsmap_local = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    from lxml import etree

    cxnSp = etree.SubElement(
        slide._element.spTree,
        qn("p:cxnSp")
    )

    # nvCxnSpPr
    nvCxnSpPr = etree.SubElement(cxnSp, qn("p:nvCxnSpPr"))
    cNvPr = etree.SubElement(nvCxnSpPr, qn("p:cNvPr"))
    cNvPr.set("id", str(100 + len(slide.shapes)))
    cNvPr.set("name", f"Connector {len(slide.shapes)}")
    cNvCxnSpPr = etree.SubElement(nvCxnSpPr, qn("p:cNvCxnSpPr"))
    nvPr = etree.SubElement(nvCxnSpPr, qn("p:nvPr"))

    # spPr
    spPr = etree.SubElement(cxnSp, qn("p:spPr"))

    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    if flipH == "1":
        xfrm.set("flipH", "1")
    if flipV == "1":
        xfrm.set("flipV", "1")

    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(int(left)))
    off.set("y", str(int(top)))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(int(width)))
    ext.set("cy", str(int(height)))

    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"))
    prstGeom.set("prst", prst)
    avLst = etree.SubElement(prstGeom, qn("a:avLst"))

    # Line properties
    ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(width_pt * 12700)))

    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", "%02X%02X%02X" % (color[0], color[1], color[2]))

    # Dash style
    if dash_style and dash_style != "solid":
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", dash_style)

    # Head end (arrow at start)
    if head_end:
        headEnd = etree.SubElement(ln, qn("a:headEnd"))
        headEnd.set("type", head_end)
        headEnd.set("w", "med")
        headEnd.set("len", "med")

    # Tail end (arrow at end)
    if tail_end:
        tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
        tailEnd.set("type", tail_end)
        tailEnd.set("w", "med")
        tailEnd.set("len", "med")

    return cxnSp


def slide1_straight_connectors(prs):
    """Slide 1: Straight connectors with different arrow styles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 1: Straight Connectors with Arrow Styles", 18)

    arrow_configs = [
        ("No arrows", None, None),
        ("Arrow at end", None, "arrow"),
        ("Arrow at start", "arrow", None),
        ("Arrows both ends", "arrow", "arrow"),
        ("Triangle end", None, "triangle"),
        ("Stealth end", None, "stealth"),
        ("Diamond end", None, "diamond"),
        ("Oval end", None, "oval"),
    ]

    for i, (label, head, tail) in enumerate(arrow_configs):
        col = i % 4
        row = i // 4
        x1 = 0.8 + (col * 3.2)
        y1 = 1.0 + (row * 3.2)
        x2 = x1 + 2.2
        y2 = y1 + 1.5

        add_endpoint_dot(slide, x1, y1)
        add_endpoint_dot(slide, x2, y2)
        add_connector(slide, x1, y1, x2, y2, "straight",
                      color=(0x44, 0x72, 0xC4), width_pt=2,
                      head_end=head, tail_end=tail)
        add_label(slide, x1, y2 + 0.2, 2.5, 0.4, label, 10)


def slide2_bent_connectors(prs):
    """Slide 2: Bent (elbow) connectors."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 2: Bent (Elbow) Connectors", 18)

    bent_configs = [
        ("Horizontal to vertical", 0.8, 1.2, 4.0, 3.0, (0x44, 0x72, 0xC4)),
        ("Vertical to horizontal", 5.0, 1.2, 8.0, 3.5, (0xED, 0x7D, 0x31)),
        ("Short bend", 9.0, 1.2, 11.5, 2.0, (0x70, 0xAD, 0x47)),
        ("Long bend down", 0.8, 4.0, 4.0, 6.5, (0xFF, 0xC0, 0x00)),
        ("Long bend right", 5.0, 4.5, 8.5, 6.5, (0x80, 0x00, 0x80)),
        ("Steep bend", 9.5, 4.0, 12.0, 7.0, (0xC0, 0x00, 0x00)),
    ]

    for label, x1, y1, x2, y2, color in bent_configs:
        add_endpoint_dot(slide, x1, y1, color)
        add_endpoint_dot(slide, x2, y2, color)
        add_connector(slide, x1, y1, x2, y2, "bent",
                      color=color, width_pt=2, tail_end="arrow")
        lx = min(x1, x2)
        ly = max(y1, y2) + 0.15
        add_label(slide, lx, ly, 3.0, 0.3, label, 9)


def slide3_curved_connectors(prs):
    """Slide 3: Curved connectors."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 3: Curved Connectors", 18)

    curve_configs = [
        ("Gentle curve", 0.8, 1.2, 4.0, 3.0, (0x44, 0x72, 0xC4)),
        ("S-curve tall", 5.0, 1.0, 8.0, 3.5, (0xED, 0x7D, 0x31)),
        ("Wide curve", 9.0, 1.2, 12.5, 2.5, (0x70, 0xAD, 0x47)),
        ("Downward curve", 0.8, 4.0, 4.0, 6.5, (0xFF, 0xC0, 0x00)),
        ("Steep curve", 5.0, 4.0, 8.0, 7.0, (0x80, 0x00, 0x80)),
        ("Short curve", 9.5, 4.5, 12.0, 6.0, (0xC0, 0x00, 0x00)),
    ]

    for label, x1, y1, x2, y2, color in curve_configs:
        add_endpoint_dot(slide, x1, y1, color)
        add_endpoint_dot(slide, x2, y2, color)
        add_connector(slide, x1, y1, x2, y2, "curved",
                      color=color, width_pt=2, tail_end="arrow")
        lx = min(x1, x2)
        ly = max(y1, y2) + 0.15
        add_label(slide, lx, ly, 3.0, 0.3, label, 9)


def slide4_line_styles(prs):
    """Slide 4: Connectors with different line styles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 4: Connector Line Styles", 18)

    line_styles = [
        ("Solid", "solid"),
        ("Dash", "dash"),
        ("Dot", "dot"),
        ("Dash-Dot", "dashDot"),
        ("Long Dash", "lgDash"),
        ("Long Dash-Dot", "lgDashDot"),
        ("Long Dash-Dot-Dot", "lgDashDotDot"),
        ("System Dash", "sysDash"),
    ]

    for i, (label, dash) in enumerate(line_styles):
        col = i % 2
        row = i // 2
        x1 = 0.8 + (col * 6.5)
        y1 = 1.0 + (row * 1.5)
        x2 = x1 + 5.0
        y2 = y1

        add_connector(slide, x1, y1, x2, y2, "straight",
                      color=(0x33, 0x33, 0x33), width_pt=2,
                      tail_end="arrow", dash_style=dash)
        add_label(slide, x1, y1 + 0.15, 5.0, 0.3, label, 11)

    # Line widths
    add_label(slide, 0.3, 7.0, 12, 0.3, "", 1)  # spacer

    widths = [
        ("0.5pt", 0.5),
        ("1pt", 1),
        ("2pt", 2),
        ("4pt", 4),
        ("6pt", 6),
    ]

    base_y = 7.2
    for i, (label, width) in enumerate(widths):
        # These would go below slide, so let's use last available rows
        pass

    # Instead, add width comparison in the remaining space
    y_start = 7.0
    if y_start < 7.5:
        pass  # Skip if no room


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide1_straight_connectors(prs)
    slide2_bent_connectors(prs)
    slide3_curved_connectors(prs)
    slide4_line_styles(prs)

    output_path = ROOT / "test-data" / "connector-stress-test.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"Generated: {output_path}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
