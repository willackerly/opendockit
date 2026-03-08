#!/usr/bin/env python3
"""
Generate a gradient stress-test PPTX that exercises gradient fills thoroughly.

Creates slides covering linear, radial, multi-stop, shape-specific, and
line/stroke gradient fills.

Usage:
    python3 scripts/generate-gradient-stress-test.py
    # Output: test-data/gradient-stress-test.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
import copy

ROOT = __import__("pathlib").Path(__file__).resolve().parent.parent

# Standard 16:9 slide dimensions
SLIDE_WIDTH = Emu(12192000)   # 13.333 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches


def add_label(slide, x, y, w, h, text, font_size=10):
    """Add a text label to the slide."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def set_linear_gradient(shape, angle_deg, stops):
    """
    Set a linear gradient fill on a shape.
    stops: list of (position_pct, RGBColor) tuples, position 0-100.
    angle_deg: rotation angle in degrees.
    """
    spPr = shape._element.spPr
    # Remove any existing fill
    for child in list(spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "gradFill", "noFill", "pattFill", "blipFill"):
            spPr.remove(child)

    # Build gradFill element
    gradFill = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = gradFill.makeelement(qn("a:gsLst"), {})

    for pos_pct, color in stops:
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(int(pos_pct * 1000))})
        srgbClr = gs.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
        gs.append(srgbClr)
        gsLst.append(gs)

    gradFill.append(gsLst)

    lin = gradFill.makeelement(qn("a:lin"), {
        "ang": str(int(angle_deg * 60000)),
        "scaled": "1"
    })
    gradFill.append(lin)

    spPr.append(gradFill)


def set_radial_gradient(shape, stops, focus_x=50, focus_y=50):
    """
    Set a radial (path) gradient fill on a shape.
    stops: list of (position_pct, (r,g,b)) tuples.
    focus_x, focus_y: center position in percent (0-100).
    """
    spPr = shape._element.spPr
    for child in list(spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "gradFill", "noFill", "pattFill", "blipFill"):
            spPr.remove(child)

    gradFill = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = gradFill.makeelement(qn("a:gsLst"), {})

    for pos_pct, color in stops:
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(int(pos_pct * 1000))})
        srgbClr = gs.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
        gs.append(srgbClr)
        gsLst.append(gs)

    gradFill.append(gsLst)

    # Path gradient (radial)
    path = gradFill.makeelement(qn("a:path"), {"path": "circle"})
    l_pct = int(focus_x * 1000)
    t_pct = int(focus_y * 1000)
    r_pct = int(focus_x * 1000)
    b_pct = int(focus_y * 1000)
    fillToRect = path.makeelement(qn("a:fillToRect"), {
        "l": str(l_pct), "t": str(t_pct),
        "r": str(r_pct), "b": str(b_pct)
    })
    path.append(fillToRect)
    gradFill.append(path)

    spPr.append(gradFill)


def set_gradient_line(shape, angle_deg, stops, width_pt=3):
    """Set a gradient stroke on a shape's outline."""
    spPr = shape._element.spPr
    # Remove existing line
    for child in list(spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "ln":
            spPr.remove(child)

    ln = spPr.makeelement(qn("a:ln"), {"w": str(int(width_pt * 12700))})
    gradFill = ln.makeelement(qn("a:gradFill"), {})
    gsLst = gradFill.makeelement(qn("a:gsLst"), {})

    for pos_pct, color in stops:
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(int(pos_pct * 1000))})
        srgbClr = gs.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
        gs.append(srgbClr)
        gsLst.append(gs)

    gradFill.append(gsLst)

    lin = gradFill.makeelement(qn("a:lin"), {
        "ang": str(int(angle_deg * 60000)),
        "scaled": "1"
    })
    gradFill.append(lin)

    ln.append(gradFill)
    spPr.append(ln)


def slide1_linear_gradients(prs):
    """Slide 1: Linear gradients at 0, 45, 90, 135, 180 degrees."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 1: Linear Gradients at Different Angles", 18)

    angles = [0, 45, 90, 135, 180]
    stops = [(0, (0x00, 0x6B, 0xBD)), (100, (0xE8, 0x3E, 0x3E))]

    for i, angle in enumerate(angles):
        x = 0.5 + (i * 2.5)
        y = 1.0
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(x), Inches(y), Inches(2.0), Inches(2.5)
        )
        set_linear_gradient(shape, angle, stops)
        add_label(slide, x, y + 2.6, 2.0, 0.4, f"{angle} degrees", 12)

    # Second row: same angles with 3 stops
    stops3 = [(0, (0xFF, 0x00, 0x00)), (50, (0xFF, 0xFF, 0x00)), (100, (0x00, 0x00, 0xFF))]
    for i, angle in enumerate(angles):
        x = 0.5 + (i * 2.5)
        y = 4.0
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(2.0), Inches(2.5)
        )
        set_linear_gradient(shape, angle, stops3)
        add_label(slide, x, y + 2.6, 2.0, 0.4, f"{angle}° (3 stops)", 10)


def slide2_radial_gradients(prs):
    """Slide 2: Radial gradients with different center positions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 2: Radial Gradients (Different Centers)", 18)

    centers = [
        (50, 50, "Center"),
        (0, 0, "Top-Left"),
        (100, 0, "Top-Right"),
        (0, 100, "Bottom-Left"),
        (100, 100, "Bottom-Right"),
    ]
    stops = [(0, (0xFF, 0xFF, 0xFF)), (100, (0x00, 0x33, 0x66))]

    for i, (fx, fy, label) in enumerate(centers):
        x = 0.5 + (i * 2.5)
        y = 1.0
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(2.0), Inches(2.5)
        )
        set_radial_gradient(shape, stops, fx, fy)
        add_label(slide, x, y + 2.6, 2.0, 0.4, label, 11)

    # Second row: radial with different color combinations
    color_combos = [
        ([(0, (0xFF, 0xD7, 0x00)), (100, (0xFF, 0x45, 0x00))], "Gold->Red"),
        ([(0, (0x00, 0xFF, 0x00)), (100, (0x00, 0x64, 0x00))], "Lime->Dark Green"),
        ([(0, (0xFF, 0xFF, 0xFF)), (50, (0x87, 0xCE, 0xEB)), (100, (0x00, 0x00, 0x8B))], "White->Sky->Navy"),
        ([(0, (0xFF, 0xC0, 0xCB)), (100, (0x80, 0x00, 0x80))], "Pink->Purple"),
        ([(0, (0xFF, 0xFF, 0x00)), (100, (0xFF, 0x69, 0xB4))], "Yellow->Hot Pink"),
    ]

    for i, (stops, label) in enumerate(color_combos):
        x = 0.5 + (i * 2.5)
        y = 4.2
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(2.0), Inches(2.5)
        )
        set_radial_gradient(shape, stops, 50, 50)
        add_label(slide, x, y + 2.6, 2.0, 0.3, label, 10)


def slide3_multi_stop_gradients(prs):
    """Slide 3: Multi-stop gradients (2, 3, 4, 5 stops)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 3: Multi-Stop Gradients", 18)

    multi_stops = [
        ("2 stops", [(0, (0x00, 0x00, 0xFF)), (100, (0xFF, 0x00, 0x00))]),
        ("3 stops", [(0, (0xFF, 0x00, 0x00)), (50, (0xFF, 0xFF, 0x00)), (100, (0x00, 0xFF, 0x00))]),
        ("4 stops", [
            (0, (0xFF, 0x00, 0x00)), (33, (0xFF, 0xFF, 0x00)),
            (66, (0x00, 0xFF, 0x00)), (100, (0x00, 0x00, 0xFF))
        ]),
        ("5 stops (rainbow)", [
            (0, (0xFF, 0x00, 0x00)), (25, (0xFF, 0xA5, 0x00)),
            (50, (0xFF, 0xFF, 0x00)), (75, (0x00, 0xFF, 0x00)),
            (100, (0x00, 0x00, 0xFF))
        ]),
    ]

    for i, (label, stops) in enumerate(multi_stops):
        x = 0.5 + (i * 3.1)
        y = 1.0
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(2.8), Inches(2.0)
        )
        set_linear_gradient(shape, 0, stops)
        add_label(slide, x, y + 2.1, 2.8, 0.3, label, 12)

    # Same stops but vertical (90 degrees)
    for i, (label, stops) in enumerate(multi_stops):
        x = 0.5 + (i * 3.1)
        y = 3.8
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(2.8), Inches(2.0)
        )
        set_linear_gradient(shape, 90, stops)
        add_label(slide, x, y + 2.1, 2.8, 0.3, f"{label} (vertical)", 10)


def slide4_gradient_shapes(prs):
    """Slide 4: Gradient fills on rectangles, ellipses, rounded rectangles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 4: Gradients on Different Shape Types", 18)

    stops = [(0, (0x66, 0x00, 0xCC)), (50, (0xCC, 0x00, 0x99)), (100, (0xFF, 0x66, 0x00))]

    # Rectangle (MSO_SHAPE.RECTANGLE = 1)
    shape1 = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(3.5), Inches(2.5))
    set_linear_gradient(shape1, 45, stops)
    add_label(slide, 0.5, 3.6, 3.5, 0.3, "Rectangle", 12)

    # Ellipse (MSO_SHAPE.OVAL = 9)
    shape2 = slide.shapes.add_shape(9, Inches(4.5), Inches(1.0), Inches(3.5), Inches(2.5))
    set_linear_gradient(shape2, 45, stops)
    add_label(slide, 4.5, 3.6, 3.5, 0.3, "Ellipse", 12)

    # Rounded Rectangle (MSO_SHAPE.ROUNDED_RECTANGLE = 5)
    shape3 = slide.shapes.add_shape(5, Inches(8.5), Inches(1.0), Inches(3.5), Inches(2.5))
    set_linear_gradient(shape3, 45, stops)
    add_label(slide, 8.5, 3.6, 3.5, 0.3, "Rounded Rectangle", 12)

    # Triangle (MSO_SHAPE.ISOCELES_TRIANGLE = 7)
    stops2 = [(0, (0x00, 0xBF, 0xFF)), (100, (0x00, 0x00, 0x80))]
    shape4 = slide.shapes.add_shape(7, Inches(0.5), Inches(4.2), Inches(3.5), Inches(2.5))
    set_linear_gradient(shape4, 90, stops2)
    add_label(slide, 0.5, 6.8, 3.5, 0.3, "Triangle", 12)

    # Diamond (MSO_SHAPE.DIAMOND = 4)
    shape5 = slide.shapes.add_shape(4, Inches(4.5), Inches(4.2), Inches(3.5), Inches(2.5))
    set_radial_gradient(shape5, stops2, 50, 50)
    add_label(slide, 4.5, 6.8, 3.5, 0.3, "Diamond (radial)", 12)

    # Pentagon (MSO_SHAPE.PENTAGON = 56)
    shape6 = slide.shapes.add_shape(56, Inches(8.5), Inches(4.2), Inches(3.5), Inches(2.5))
    set_linear_gradient(shape6, 135, stops)
    add_label(slide, 8.5, 6.8, 3.5, 0.3, "Pentagon", 12)


def slide5_gradient_lines(prs):
    """Slide 5: Gradient line/stroke fills."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 5: Gradient Stroke / Line Fills", 18)

    # Shapes with gradient outlines and no fill
    line_configs = [
        ("Thin gradient line (2pt)", 2,
         [(0, (0xFF, 0x00, 0x00)), (100, (0x00, 0x00, 0xFF))], 0),
        ("Medium gradient line (4pt)", 4,
         [(0, (0x00, 0xFF, 0x00)), (50, (0xFF, 0xFF, 0x00)), (100, (0xFF, 0x00, 0x00))], 45),
        ("Thick gradient line (8pt)", 8,
         [(0, (0xFF, 0x69, 0xB4)), (100, (0x80, 0x00, 0x80))], 90),
    ]

    for i, (label, width, stops, angle) in enumerate(line_configs):
        x = 0.5 + (i * 4.2)
        y = 1.0
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.5), Inches(2.0))
        # Set no fill
        spPr = shape._element.spPr
        for child in list(spPr):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag in ("solidFill", "gradFill", "pattFill", "blipFill"):
                spPr.remove(child)
        noFill = spPr.makeelement(qn("a:noFill"), {})
        spPr.insert(0, noFill)

        set_gradient_line(shape, angle, stops, width)
        add_label(slide, x, y + 2.1, 3.5, 0.4, label, 10)

    # Shapes with both gradient fill and gradient line
    combo_configs = [
        ("Fill + line gradient", 3,
         [(0, (0xAD, 0xD8, 0xE6)), (100, (0x00, 0x00, 0x8B))],
         [(0, (0xFF, 0xD7, 0x00)), (100, (0xFF, 0x00, 0x00))]),
        ("Radial fill + line gradient", 3,
         [(0, (0xFF, 0xFF, 0xFF)), (100, (0x8B, 0x00, 0x8B))],
         [(0, (0x00, 0xFF, 0x00)), (100, (0x00, 0x00, 0x80))]),
    ]

    for i, (label, width, fill_stops, line_stops) in enumerate(combo_configs):
        x = 0.5 + (i * 6.2)
        y = 4.0
        shape = slide.shapes.add_shape(
            5, Inches(x), Inches(y), Inches(5.5), Inches(2.5)
        )
        if i == 0:
            set_linear_gradient(shape, 45, fill_stops)
        else:
            set_radial_gradient(shape, fill_stops, 50, 50)
        set_gradient_line(shape, 0, line_stops, width)
        add_label(slide, x, y + 2.6, 5.5, 0.4, label, 11)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide1_linear_gradients(prs)
    slide2_radial_gradients(prs)
    slide3_multi_stop_gradients(prs)
    slide4_gradient_shapes(prs)
    slide5_gradient_lines(prs)

    output_path = ROOT / "test-data" / "gradient-stress-test.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"Generated: {output_path}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
