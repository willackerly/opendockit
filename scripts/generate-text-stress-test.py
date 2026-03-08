#!/usr/bin/env python3
"""
Generate a text rendering stress-test PPTX.

Creates slides covering text alignment, bullet types, autofit, text body
rotation, character spacing, and paragraph spacing.

Usage:
    python3 scripts/generate-text-stress-test.py
    # Output: test-data/text-stress-test.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

ROOT = __import__("pathlib").Path(__file__).resolve().parent.parent

SLIDE_WIDTH = Emu(12192000)   # 13.333 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches

LOREM = (
    "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do eiusmod "
    "tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, "
    "quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat."
)

LOREM_SHORT = "The quick brown fox jumps over the lazy dog. Pack my box with five dozen liquor jugs."


def add_label(slide, x, y, w, h, text, font_size=10, color=None):
    """Add a text label to the slide."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = "Calibri"
    run.font.color.rgb = color or RGBColor(0x33, 0x33, 0x33)
    return txBox


def set_shape_fill(shape, color):
    """Set solid fill on a shape."""
    spPr = shape._element.spPr
    for child in list(spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "gradFill", "noFill", "pattFill", "blipFill"):
            spPr.remove(child)
    solidFill = spPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {
        "val": "%02X%02X%02X" % (color[0], color[1], color[2])
    })
    solidFill.append(srgbClr)
    spPr.append(solidFill)


def set_shape_outline(shape, color, width_pt=1):
    """Set outline on a shape."""
    spPr = shape._element.spPr
    for child in list(spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "ln":
            spPr.remove(child)
    ln = spPr.makeelement(qn("a:ln"), {"w": str(int(width_pt * 12700))})
    solidFill = ln.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {
        "val": "%02X%02X%02X" % (color[0], color[1], color[2])
    })
    solidFill.append(srgbClr)
    ln.append(solidFill)
    spPr.append(ln)


def slide1_text_alignments(prs):
    """Slide 1: All 5 text alignments."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 1: Text Alignments (Left, Center, Right, Justify, Distributed)", 18)

    alignments = [
        ("Left Aligned", PP_ALIGN.LEFT),
        ("Center Aligned", PP_ALIGN.CENTER),
        ("Right Aligned", PP_ALIGN.RIGHT),
        ("Justified Text", PP_ALIGN.JUSTIFY),
        ("Distributed", PP_ALIGN.DISTRIBUTE),
    ]

    for i, (label, align) in enumerate(alignments):
        y = 0.8 + (i * 1.3)

        # Label
        add_label(slide, 0.3, y, 2.5, 0.3, label, 11, RGBColor(0x00, 0x52, 0x8A))

        # Text box with alignment
        shape = slide.shapes.add_textbox(Inches(0.3), Inches(y + 0.35), Inches(12.5), Inches(0.8))
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = LOREM_SHORT
        run.font.size = Pt(12)
        run.font.name = "Calibri"

        # Add border for visibility
        set_shape_outline(shape, (0xCC, 0xCC, 0xCC), 0.5)


def slide2_bullet_types(prs):
    """Slide 2: Bullet types (numbered, lettered, symbol, multi-level)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 2: Bullet Types and Numbering", 18)

    # Column 1: Symbol bullets
    shape1 = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), Inches(3.8), Inches(3.0))
    tf = shape1.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Symbol Bullets:"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = "Calibri"

    bullets = [
        ("\u2022 Standard bullet point", False),
        ("\u25CB Circle bullet point", False),
        ("\u25A0 Square bullet point", False),
        ("\u2713 Checkmark bullet", False),
        ("\u2192 Arrow bullet point", False),
        ("\u2605 Star bullet point", False),
    ]
    for text, bold in bullets:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size = Pt(11)
        r.font.name = "Calibri"
        r.font.bold = bold

    set_shape_outline(shape1, (0xDD, 0xDD, 0xDD), 0.5)

    # Column 2: Numbered lists
    shape2 = slide.shapes.add_textbox(Inches(4.5), Inches(0.8), Inches(3.8), Inches(3.0))
    tf2 = shape2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    r = p.add_run()
    r.text = "Numbered Lists:"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = "Calibri"

    for i in range(1, 7):
        p = tf2.add_paragraph()
        r = p.add_run()
        r.text = f"{i}. Numbered item {i}"
        r.font.size = Pt(11)
        r.font.name = "Calibri"

    set_shape_outline(shape2, (0xDD, 0xDD, 0xDD), 0.5)

    # Column 3: Lettered lists
    shape3 = slide.shapes.add_textbox(Inches(8.7), Inches(0.8), Inches(3.8), Inches(3.0))
    tf3 = shape3.text_frame
    tf3.word_wrap = True

    p = tf3.paragraphs[0]
    r = p.add_run()
    r.text = "Lettered Lists:"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = "Calibri"

    letters = "abcdef"
    for i, letter in enumerate(letters):
        p = tf3.add_paragraph()
        r = p.add_run()
        r.text = f"{letter}) Lettered item {letter.upper()}"
        r.font.size = Pt(11)
        r.font.name = "Calibri"

    set_shape_outline(shape3, (0xDD, 0xDD, 0xDD), 0.5)

    # Multi-level bullets
    shape4 = slide.shapes.add_textbox(Inches(0.3), Inches(4.2), Inches(12.2), Inches(2.8))
    tf4 = shape4.text_frame
    tf4.word_wrap = True

    p = tf4.paragraphs[0]
    r = p.add_run()
    r.text = "Multi-Level Outline:"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = "Calibri"

    multi_level = [
        ("1. First level item", 0, True),
        ("   a. Second level sub-item", 1, False),
        ("   b. Second level sub-item", 1, False),
        ("      i. Third level detail", 2, False),
        ("      ii. Third level detail", 2, False),
        ("   c. Second level sub-item", 1, False),
        ("2. First level item", 0, True),
        ("   a. Second level sub-item", 1, False),
        ("3. First level item", 0, True),
    ]

    for text, level, bold in multi_level:
        p = tf4.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size = Pt(11)
        r.font.name = "Calibri"
        r.font.bold = bold

    set_shape_outline(shape4, (0xDD, 0xDD, 0xDD), 0.5)


def slide3_autofit(prs):
    """Slide 3: Text autofit variations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 3: Text AutoFit Variations", 18)

    long_text = LOREM + " " + LOREM

    # No autofit (overflow)
    shape1 = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), Inches(3.8), Inches(2.5))
    tf1 = shape1.text_frame
    tf1.word_wrap = True
    tf1.auto_size = MSO_AUTO_SIZE.NONE
    p = tf1.paragraphs[0]
    r = p.add_run()
    r.text = long_text
    r.font.size = Pt(14)
    r.font.name = "Calibri"
    set_shape_outline(shape1, (0xFF, 0x00, 0x00), 1)
    add_label(slide, 0.3, 3.5, 3.8, 0.3, "No AutoFit (may overflow)", 10, RGBColor(0xFF, 0x00, 0x00))

    # Shrink text on overflow
    shape2 = slide.shapes.add_textbox(Inches(4.5), Inches(0.8), Inches(3.8), Inches(2.5))
    tf2 = shape2.text_frame
    tf2.word_wrap = True
    tf2.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf2.paragraphs[0]
    r = p.add_run()
    r.text = long_text
    r.font.size = Pt(14)
    r.font.name = "Calibri"
    set_shape_outline(shape2, (0x00, 0x80, 0x00), 1)
    add_label(slide, 4.5, 3.5, 3.8, 0.3, "Shrink Text to Fit", 10, RGBColor(0x00, 0x80, 0x00))

    # Shape grows to fit text
    shape3 = slide.shapes.add_textbox(Inches(8.7), Inches(0.8), Inches(3.8), Inches(1.0))
    tf3 = shape3.text_frame
    tf3.word_wrap = True
    tf3.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = tf3.paragraphs[0]
    r = p.add_run()
    r.text = long_text
    r.font.size = Pt(14)
    r.font.name = "Calibri"
    set_shape_outline(shape3, (0x00, 0x00, 0xFF), 1)
    add_label(slide, 8.7, 3.5, 3.8, 0.3, "Shape Grows to Fit Text", 10, RGBColor(0x00, 0x00, 0xFF))

    # Small boxes with different amounts of text
    amounts = [
        ("Short", "Hello world"),
        ("Medium", "The quick brown fox jumps over the lazy dog."),
        ("Long", LOREM),
        ("Very Long", long_text),
    ]

    for i, (label, text) in enumerate(amounts):
        x = 0.3 + (i * 3.2)
        y = 4.2
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.8), Inches(1.5))
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(14)
        r.font.name = "Calibri"
        set_shape_outline(shape, (0x66, 0x66, 0x66), 0.5)
        add_label(slide, x, y + 1.6, 2.8, 0.3, f"Shrink: {label}", 9)


def slide4_text_rotation(prs):
    """Slide 4: Text body rotation (0, 90, 270)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 4: Text Body Rotation", 18)

    rotations = [
        ("No rotation (0\u00b0)", None),
        ("Rotated 90\u00b0 CW (vert)", "vert"),
        ("Rotated 270\u00b0 CW (vert270)", "vert270"),
        ("Word Art Vertical", "wordArtVert"),
    ]

    for i, (label, rot) in enumerate(rotations):
        x = 0.5 + (i * 3.2)
        y = 0.8

        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.8), Inches(5.0))
        set_shape_fill(shape, (0xE8, 0xF0, 0xFE))
        set_shape_outline(shape, (0x44, 0x72, 0xC4), 1)

        tf = shape.text_frame
        tf.word_wrap = True
        if rot:
            # Set text rotation via XML
            bodyPr = tf._txBody.find(qn("a:bodyPr"))
            if bodyPr is not None:
                bodyPr.set("vert", rot)

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "Rotated Text Sample"
        r.font.size = Pt(16)
        r.font.name = "Calibri"
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = LOREM_SHORT
        r2.font.size = Pt(11)
        r2.font.name = "Calibri"

        add_label(slide, x, y + 5.2, 2.8, 0.5, label, 10)


def slide5_character_formatting(prs):
    """Slide 5: Character spacing, bold, italic, underline combinations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 5: Character Formatting Combinations", 18)

    # Character spacing variations
    add_label(slide, 0.3, 0.7, 12, 0.3, "Character Spacing:", 12, RGBColor(0x00, 0x52, 0x8A))

    spacings = [
        ("Very Tight (-3pt)", -300),
        ("Tight (-1.5pt)", -150),
        ("Normal (0pt)", 0),
        ("Loose (+1.5pt)", 150),
        ("Very Loose (+3pt)", 300),
    ]

    for i, (label, spacing) in enumerate(spacings):
        y = 1.1 + (i * 0.4)
        shape = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12), Inches(0.35))
        tf = shape.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{label}: ABCDEFGHIJKLMNOP abcdefghijklmnop"
        r.font.size = Pt(14)
        r.font.name = "Calibri"
        # Set character spacing via XML
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(spacing))

    # Bold/Italic/Underline combinations
    add_label(slide, 0.3, 3.3, 12, 0.3, "Formatting Combinations:", 12, RGBColor(0x00, 0x52, 0x8A))

    combos = [
        ("Normal text", False, False, False),
        ("Bold text", True, False, False),
        ("Italic text", False, True, False),
        ("Underline text", False, False, True),
        ("Bold + Italic", True, True, False),
        ("Bold + Underline", True, False, True),
        ("Italic + Underline", False, True, True),
        ("Bold + Italic + Underline", True, True, True),
    ]

    for i, (text, bold, italic, underline) in enumerate(combos):
        col = i % 4
        row = i // 4
        x = 0.5 + (col * 3.1)
        y = 3.7 + (row * 0.5)

        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.8), Inches(0.4))
        tf = shape.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(13)
        r.font.name = "Calibri"
        r.font.bold = bold
        r.font.italic = italic
        r.font.underline = underline

    # Strikethrough and superscript/subscript
    add_label(slide, 0.3, 5.0, 12, 0.3, "Special Formatting:", 12, RGBColor(0x00, 0x52, 0x8A))

    shape_strike = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(4), Inches(0.4))
    tf = shape_strike.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Strikethrough text"
    r.font.size = Pt(14)
    r.font.name = "Calibri"
    # Strikethrough via XML
    rPr = r._r.get_or_add_rPr()
    rPr.set("strike", "sngStrike")

    # Superscript / subscript
    shape_super = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(10), Inches(0.5))
    tf = shape_super.text_frame
    p = tf.paragraphs[0]

    r1 = p.add_run()
    r1.text = "E = mc"
    r1.font.size = Pt(16)
    r1.font.name = "Calibri"

    r2 = p.add_run()
    r2.text = "2"
    r2.font.size = Pt(16)
    r2.font.name = "Calibri"
    rPr2 = r2._r.get_or_add_rPr()
    rPr2.set("baseline", "30000")  # superscript

    r3 = p.add_run()
    r3.text = "    H"
    r3.font.size = Pt(16)
    r3.font.name = "Calibri"

    r4 = p.add_run()
    r4.text = "2"
    r4.font.size = Pt(16)
    r4.font.name = "Calibri"
    rPr4 = r4._r.get_or_add_rPr()
    rPr4.set("baseline", "-25000")  # subscript

    r5 = p.add_run()
    r5.text = "O    x"
    r5.font.size = Pt(16)
    r5.font.name = "Calibri"

    r6 = p.add_run()
    r6.text = "n"
    r6.font.size = Pt(16)
    r6.font.name = "Calibri"
    rPr6 = r6._r.get_or_add_rPr()
    rPr6.set("baseline", "30000")

    r7 = p.add_run()
    r7.text = " + y"
    r7.font.size = Pt(16)
    r7.font.name = "Calibri"

    r8 = p.add_run()
    r8.text = "n"
    r8.font.size = Pt(16)
    r8.font.name = "Calibri"
    rPr8 = r8._r.get_or_add_rPr()
    rPr8.set("baseline", "30000")

    r9 = p.add_run()
    r9.text = " = z"
    r9.font.size = Pt(16)
    r9.font.name = "Calibri"

    r10 = p.add_run()
    r10.text = "n"
    r10.font.size = Pt(16)
    r10.font.name = "Calibri"
    rPr10 = r10._r.get_or_add_rPr()
    rPr10.set("baseline", "30000")

    add_label(slide, 0.5, 6.5, 10, 0.3, "Superscript (E=mc\u00b2) and Subscript (H\u2082O)", 10)


def slide6_paragraph_spacing(prs):
    """Slide 6: Multiple paragraphs with different spacing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 6: Paragraph Spacing (Before, After, Line Spacing)", 18)

    # Column 1: Space before/after
    add_label(slide, 0.3, 0.7, 3.8, 0.3, "Space Before/After:", 11, RGBColor(0x00, 0x52, 0x8A))

    shape1 = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(3.8), Inches(5.5))
    tf1 = shape1.text_frame
    tf1.word_wrap = True
    set_shape_outline(shape1, (0xCC, 0xCC, 0xCC), 0.5)

    space_configs = [
        ("No extra spacing", 0, 0),
        ("6pt before", 6, 0),
        ("6pt after", 0, 6),
        ("12pt before", 12, 0),
        ("12pt after", 0, 12),
        ("12pt before + 6pt after", 12, 6),
    ]

    for i, (text, before, after) in enumerate(space_configs):
        if i == 0:
            p = tf1.paragraphs[0]
        else:
            p = tf1.add_paragraph()
        r = p.add_run()
        r.text = text
        r.font.size = Pt(12)
        r.font.name = "Calibri"

        pPr = p._p.get_or_add_pPr()
        spcBef = pPr.makeelement(qn("a:spcBef"), {})
        spcPts_b = spcBef.makeelement(qn("a:spcPts"), {"val": str(before * 100)})
        spcBef.append(spcPts_b)
        pPr.append(spcBef)

        spcAft = pPr.makeelement(qn("a:spcAft"), {})
        spcPts_a = spcAft.makeelement(qn("a:spcPts"), {"val": str(after * 100)})
        spcAft.append(spcPts_a)
        pPr.append(spcAft)

    # Column 2: Line spacing
    add_label(slide, 4.5, 0.7, 3.8, 0.3, "Line Spacing:", 11, RGBColor(0x00, 0x52, 0x8A))

    line_spacings = [
        ("Single (100%)", 100),
        ("1.15x (115%)", 115),
        ("1.5x (150%)", 150),
        ("Double (200%)", 200),
        ("Triple (300%)", 300),
    ]

    y_offset = 1.1
    for label, pct in line_spacings:
        shape = slide.shapes.add_textbox(Inches(4.5), Inches(y_offset), Inches(3.8), Inches(1.0))
        tf = shape.text_frame
        tf.word_wrap = True
        set_shape_outline(shape, (0xCC, 0xCC, 0xCC), 0.5)

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{label}: {LOREM_SHORT[:60]}"
        r.font.size = Pt(10)
        r.font.name = "Calibri"

        # Set line spacing
        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
        spcPct = lnSpc.makeelement(qn("a:spcPct"), {"val": str(pct * 1000)})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)

        y_offset += 1.2

    # Column 3: Fixed line spacing (points)
    add_label(slide, 8.7, 0.7, 3.8, 0.3, "Fixed Line Spacing (points):", 11, RGBColor(0x00, 0x52, 0x8A))

    fixed_spacings = [
        ("12pt line spacing", 12),
        ("16pt line spacing", 16),
        ("24pt line spacing", 24),
        ("36pt line spacing", 36),
    ]

    y_offset = 1.1
    for label, pts in fixed_spacings:
        shape = slide.shapes.add_textbox(Inches(8.7), Inches(y_offset), Inches(3.8), Inches(1.2))
        tf = shape.text_frame
        tf.word_wrap = True
        set_shape_outline(shape, (0xCC, 0xCC, 0xCC), 0.5)

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{label}: {LOREM_SHORT[:50]}"
        r.font.size = Pt(10)
        r.font.name = "Calibri"

        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
        spcPts = lnSpc.makeelement(qn("a:spcPts"), {"val": str(pts * 100)})
        lnSpc.append(spcPts)
        pPr.append(lnSpc)

        y_offset += 1.5


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide1_text_alignments(prs)
    slide2_bullet_types(prs)
    slide3_autofit(prs)
    slide4_text_rotation(prs)
    slide5_character_formatting(prs)
    slide6_paragraph_spacing(prs)

    output_path = ROOT / "test-data" / "text-stress-test.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"Generated: {output_path}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
