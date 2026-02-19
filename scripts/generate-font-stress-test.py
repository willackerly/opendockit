#!/usr/bin/env python3
"""
Generate a font stress-test PPTX that exercises every bundled font family.

Creates slides with sample text in all 42 bundled families, including
bold/italic variants, different sizes, and mixed-font paragraphs.

Usage:
    python3 scripts/generate-font-stress-test.py
    # Output: test-data/font-stress-test.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# All 42 bundled font families grouped by category
OFFICE_SUBSTITUTES = [
    ("Calibri", "Carlito substitute"),
    ("Calibri Light", "Carlito Light substitute"),
    ("Cambria", "Caladea substitute"),
    ("Arial", "Liberation Sans substitute"),
    ("Arial Narrow", "Liberation Sans Narrow substitute"),
    ("Times New Roman", "Liberation Serif substitute"),
    ("Courier New", "Liberation Mono substitute"),
    ("Segoe UI", "Selawik substitute"),
    ("Segoe UI Light", "Selawik Light substitute"),
    ("Segoe UI Semibold", "Selawik Semibold substitute"),
    ("Segoe UI Semilight", "Selawik Semilight substitute"),
    ("Georgia", "Gelasio substitute"),
    ("Palatino Linotype", "TeX Gyre Pagella substitute"),
    ("Bookman Old Style", "TeX Gyre Bonum substitute"),
    ("Century Schoolbook", "TeX Gyre Schola substitute"),
]

GOOGLE_FONTS_SANS = [
    ("Arimo", "sans-serif"),
    ("Barlow", "sans-serif"),
    ("Barlow Light", "sans-serif, light weight"),
    ("Comfortaa", "rounded sans-serif"),
    ("Lato", "sans-serif"),
    ("Lato Light", "sans-serif, light weight"),
    ("Montserrat", "geometric sans-serif"),
    ("Noto Sans", "sans-serif"),
    ("Open Sans", "humanist sans-serif"),
    ("Oswald", "narrow sans-serif"),
    ("Poppins", "geometric sans-serif"),
    ("Raleway", "elegant sans-serif"),
    ("Roboto", "neo-grotesque sans-serif"),
    ("Source Sans Pro", "sans-serif"),
    ("Ubuntu", "humanist sans-serif"),
    ("Play", "sans-serif"),
]

GOOGLE_FONTS_SERIF = [
    ("Noto Serif", "serif"),
    ("Playfair Display", "transitional serif"),
    ("Roboto Slab", "slab serif"),
    ("Roboto Slab Light", "slab serif, light weight"),
    ("Roboto Slab SemiBold", "slab serif, semibold weight"),
    ("Tinos", "transitional serif"),
]

GOOGLE_FONTS_MONO = [
    ("Courier Prime", "monospace"),
    ("Fira Code", "monospace"),
    ("Roboto Mono", "monospace"),
    ("Source Code Pro", "monospace"),
]

GOOGLE_FONTS_SYMBOL = [
    ("Noto Sans Symbols", "symbols & arrows"),
]

SLIDE_WIDTH = Emu(12192000)   # 10 inches (standard widescreen)
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches

SAMPLE_TEXT = "The quick brown fox jumps over the lazy dog. 0123456789"
SAMPLE_SHORT = "AaBbCcDdEeFfGg 0123456789"
SYMBOL_TEXT = "\u2190\u2191\u2192\u2193 \u2600\u2601\u2602 \u25A0\u25B2\u25CF \u2200\u2203\u2208 \u00A9\u00AE\u2122"


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(18)
    run2.font.name = "Calibri"
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    return slide


def add_font_showcase_slide(prs, fonts, category_title, bg_color=None):
    """Add a slide showcasing multiple fonts with sample text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    if bg_color:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Category header
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = category_title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x99)

    y_pos = 0.6
    line_height = 0.42  # tighter spacing

    for font_name, description in fonts:
        if y_pos > 7.0:
            break

        # Font name label
        txBox = slide.shapes.add_textbox(
            Inches(0.3), Inches(y_pos), Inches(9.4), Inches(line_height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        # Label in small Calibri
        label_run = p.add_run()
        label_run.text = f"{font_name} ({description}): "
        label_run.font.size = Pt(8)
        label_run.font.name = "Calibri"
        label_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        # Sample in the actual font - regular
        sample_run = p.add_run()
        sample_run.text = SAMPLE_SHORT
        sample_run.font.size = Pt(14)
        sample_run.font.name = font_name
        sample_run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

        y_pos += line_height

        # Bold + Italic row
        txBox2 = slide.shapes.add_textbox(
            Inches(0.3), Inches(y_pos), Inches(9.4), Inches(line_height)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]

        # Bold sample
        bold_run = p2.add_run()
        bold_run.text = "Bold "
        bold_run.font.size = Pt(12)
        bold_run.font.name = font_name
        bold_run.font.bold = True
        bold_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Italic sample
        italic_run = p2.add_run()
        italic_run.text = "Italic "
        italic_run.font.size = Pt(12)
        italic_run.font.name = font_name
        italic_run.font.italic = True
        italic_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Bold Italic sample
        bi_run = p2.add_run()
        bi_run.text = "Bold Italic"
        bi_run.font.size = Pt(12)
        bi_run.font.name = font_name
        bi_run.font.bold = True
        bi_run.font.italic = True
        bi_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        y_pos += line_height + 0.08  # extra gap between families

    return slide


def add_size_comparison_slide(prs, fonts, title):
    """Add a slide showing the same text at multiple sizes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x99)

    y_pos = 0.6
    sizes = [10, 14, 20, 28, 36]

    for font_name, _ in fonts[:5]:  # Top 5 fonts
        for size in sizes:
            if y_pos > 7.2:
                break
            txBox = slide.shapes.add_textbox(
                Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]

            label = p.add_run()
            label.text = f"{font_name} {size}pt: "
            label.font.size = Pt(8)
            label.font.name = "Calibri"
            label.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

            sample = p.add_run()
            sample.text = "Hamburgevons"
            sample.font.size = Pt(size)
            sample.font.name = font_name
            sample.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

            y_pos += max(0.35, size / 72 * 1.4)

        y_pos += 0.15

    return slide


def add_mixed_paragraph_slide(prs):
    """Add a slide with paragraphs mixing multiple fonts (real-world scenario)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Mixed Font Paragraphs (Real-World Scenarios)"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x99)

    # Mixed paragraph 1: presentation body text
    txBox1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(1.5))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]

    pairs = [
        ("Revenue grew ", "Calibri", False, False),
        ("23.5%", "Calibri", True, False),
        (" year-over-year, driven by ", "Calibri", False, False),
        ("strong Q4 performance", "Calibri", True, True),
        (". Source: ", "Calibri", False, False),
        ("Bloomberg Terminal", "Courier New", False, False),
        (" data feed.", "Calibri", False, False),
    ]
    for text, font, bold, italic in pairs:
        r = p1.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(16)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    # Mixed paragraph 2: code in presentation
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]

    code_pairs = [
        ("The function ", "Segoe UI", False),
        ("calculateMetrics()", "Fira Code", False),
        (" accepts a ", "Segoe UI", False),
        ("FontFamily", "Source Code Pro", True),
        (" parameter and returns ", "Segoe UI", False),
        ("AdvanceWidth[]", "Roboto Mono", True),
        (".", "Segoe UI", False),
    ]
    for text, font, bold in code_pairs:
        r = p2.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(14)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)

    # Mixed paragraph 3: multilingual / symbols
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1.0))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]

    sym_pairs = [
        ("Arrows: ", "Noto Sans", False),
        ("\u2190 \u2191 \u2192 \u2193 \u21D0 \u21D2", "Noto Sans Symbols", False),
        ("  Math: ", "Noto Sans", False),
        ("\u2200x \u2208 \u211D, \u2203y: x\u00B2 + y\u00B2 = 1", "Noto Sans Symbols", False),
        ("  Currency: \u00A3 \u20AC \u00A5 \u20B9", "Noto Sans", False),
    ]
    for text, font, bold in sym_pairs:
        r = p3.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(16)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    # Heading hierarchy
    y = 4.9
    heading_fonts = [
        ("Heading 1 \u2014 Montserrat Bold", "Montserrat", 28, True),
        ("Heading 2 \u2014 Raleway Regular", "Raleway", 22, False),
        ("Heading 3 \u2014 Poppins Bold", "Poppins", 18, True),
        ("Body text in Roboto. Lorem ipsum dolor sit amet, consectetur adipiscing elit.", "Roboto", 13, False),
        ("Caption in Lato Light italic", "Lato Light", 11, False),
    ]
    for text, font, size, bold in heading_fonts:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        y += max(0.35, size / 72 * 1.5)

    return slide


def add_full_alphabet_slide(prs, font_name, description):
    """Add a slide with full character set for a single font."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{font_name} \u2014 {description}"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x99)

    lines = [
        ("Regular 24pt", font_name, 24, False, False),
        ("ABCDEFGHIJKLMNOPQRSTUVWXYZ", font_name, 24, False, False),
        ("abcdefghijklmnopqrstuvwxyz", font_name, 24, False, False),
        ("0123456789 !@#$%^&*()+-=[]{}|;':\",./<>?", font_name, 20, False, False),
        ("", font_name, 8, False, False),
        ("Bold 20pt", font_name, 20, True, False),
        ("The quick brown fox jumps over the lazy dog.", font_name, 20, True, False),
        ("", font_name, 8, False, False),
        ("Italic 20pt", font_name, 20, False, True),
        ("The quick brown fox jumps over the lazy dog.", font_name, 20, False, True),
        ("", font_name, 8, False, False),
        ("Bold Italic 18pt", font_name, 18, True, True),
        ("The quick brown fox jumps over the lazy dog.", font_name, 18, True, True),
        ("", font_name, 8, False, False),
        ("Small 10pt", font_name, 10, False, False),
        ("Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris.", font_name, 10, False, False),
    ]

    y = 0.55
    for text, font, size, bold, italic in lines:
        if y > 7.2:
            break
        h = max(0.3, size / 72 * 1.6)
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(y), Inches(9.4), Inches(h + 0.1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        if text.endswith("pt"):
            r.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            r.font.name = "Calibri"
            r.font.size = Pt(9)
            h = 0.2
        else:
            r.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        y += h

    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    add_title_slide(
        prs,
        "OpenDocKit Font Stress Test",
        "42 bundled font families \u2022 130 faces \u2022 5.0 MB WOFF2 \u2022 100% offline"
    )

    # Slide 2-3: Office font substitutes showcase
    add_font_showcase_slide(prs, OFFICE_SUBSTITUTES[:8], "Office Core Font Substitutes (1/2)")
    add_font_showcase_slide(prs, OFFICE_SUBSTITUTES[8:], "Office Core Font Substitutes (2/2)")

    # Slide 4-5: Google Sans fonts
    add_font_showcase_slide(prs, GOOGLE_FONTS_SANS[:8], "Google Fonts \u2014 Sans-Serif (1/2)")
    add_font_showcase_slide(prs, GOOGLE_FONTS_SANS[8:], "Google Fonts \u2014 Sans-Serif (2/2)")

    # Slide 6: Google Serif + Mono + Symbol
    add_font_showcase_slide(prs, GOOGLE_FONTS_SERIF, "Google Fonts \u2014 Serif")
    add_font_showcase_slide(
        prs,
        GOOGLE_FONTS_MONO + GOOGLE_FONTS_SYMBOL,
        "Google Fonts \u2014 Monospace & Symbols"
    )

    # Slide 8: Size comparison
    add_size_comparison_slide(
        prs,
        [("Calibri", ""), ("Arial", ""), ("Times New Roman", ""), ("Roboto", ""), ("Montserrat", "")],
        "Size Comparison \u2014 Core Fonts at Multiple Sizes"
    )

    # Slide 9: Mixed paragraphs
    add_mixed_paragraph_slide(prs)

    # Slides 10+: Full alphabet for key fonts
    key_fonts = [
        ("Calibri", "Office default body font (Carlito substitute)"),
        ("Arial", "Universal sans-serif (Liberation Sans)"),
        ("Times New Roman", "Classic serif (Liberation Serif)"),
        ("Courier New", "Classic monospace (Liberation Mono)"),
        ("Georgia", "Screen-optimized serif (Gelasio)"),
        ("Segoe UI", "Windows UI font (Selawik)"),
        ("Roboto", "Google's flagship sans-serif"),
        ("Montserrat", "Geometric display sans-serif"),
        ("Playfair Display", "High-contrast serif display"),
        ("Fira Code", "Programming ligature monospace"),
        ("Noto Sans Symbols", "Unicode symbols & arrows"),
    ]
    for font_name, desc in key_fonts:
        add_full_alphabet_slide(prs, font_name, desc)

    output_path = ROOT / "test-data" / "font-stress-test.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    slide_count = len(prs.slides)
    font_count = len(OFFICE_SUBSTITUTES) + len(GOOGLE_FONTS_SANS) + len(GOOGLE_FONTS_SERIF) + len(GOOGLE_FONTS_MONO) + len(GOOGLE_FONTS_SYMBOL)
    print(f"Generated: {output_path}")
    print(f"  Slides: {slide_count}")
    print(f"  Font families exercised: {font_count}")


ROOT = __import__("pathlib").Path(__file__).resolve().parent.parent

if __name__ == "__main__":
    main()
