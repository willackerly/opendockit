#!/usr/bin/env python3
"""
Generate an effect stress-test PPTX that exercises visual effects thoroughly.

Creates slides covering drop shadows, outer glow, reflection, soft edges,
and combined effects.

Usage:
    python3 scripts/generate-effect-stress-test.py
    # Output: test-data/effect-stress-test.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

ROOT = __import__("pathlib").Path(__file__).resolve().parent.parent

SLIDE_WIDTH = Emu(12192000)   # 13.333 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches


def add_label(slide, x, y, w, h, text, font_size=10):
    """Add a text label to the slide."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def set_solid_fill(shape, color):
    """Set a solid fill on a shape."""
    spPr = shape._element.spPr
    for child in list(spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "gradFill", "noFill", "pattFill", "blipFill"):
            spPr.remove(child)
    solidFill = spPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {
        "val": "%02X%02X%02X" % (color[0], color[1], color[2])
    })
    solidFill.append(srgbClr)
    spPr.append(solidFill)


def add_drop_shadow(shape, dist_pt, blur_pt, angle_deg, color, alpha_pct=50):
    """
    Add a drop shadow effect to a shape.
    dist_pt: shadow distance in points
    blur_pt: blur radius in points
    angle_deg: direction angle in degrees
    color: (r, g, b) tuple
    alpha_pct: opacity 0-100
    """
    spPr = shape._element.spPr

    # Get or create effectLst
    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(effectLst)

    outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur_pt * 12700)),
        "dist": str(int(dist_pt * 12700)),
        "dir": str(int(angle_deg * 60000)),
        "algn": "bl",
        "rotWithShape": "0",
    })

    srgbClr = outerShdw.makeelement(qn("a:srgbClr"), {
        "val": "%02X%02X%02X" % (color[0], color[1], color[2])
    })
    alpha = srgbClr.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))})
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)


def add_outer_glow(shape, blur_pt, color, alpha_pct=40):
    """Add an outer glow effect to a shape."""
    spPr = shape._element.spPr

    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(effectLst)

    glow = effectLst.makeelement(qn("a:glow"), {
        "rad": str(int(blur_pt * 12700)),
    })

    srgbClr = glow.makeelement(qn("a:srgbClr"), {
        "val": "%02X%02X%02X" % (color[0], color[1], color[2])
    })
    alpha = srgbClr.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))})
    srgbClr.append(alpha)
    glow.append(srgbClr)
    effectLst.append(glow)


def add_reflection(shape, blur_pt=1, start_alpha_pct=50, end_pos_pct=50, dist_pt=0, dir_deg=90, fade_dir_deg=90):
    """Add a reflection effect to a shape."""
    spPr = shape._element.spPr

    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(effectLst)

    reflection = effectLst.makeelement(qn("a:reflection"), {
        "blurRad": str(int(blur_pt * 12700)),
        "stA": str(int(start_alpha_pct * 1000)),
        "endA": "0",
        "endPos": str(int(end_pos_pct * 1000)),
        "dist": str(int(dist_pt * 12700)),
        "dir": str(int(dir_deg * 60000)),
        "fadeDir": str(int(fade_dir_deg * 60000)),
        "algn": "bl",
        "rotWithShape": "0",
    })
    effectLst.append(reflection)


def add_soft_edge(shape, radius_pt):
    """Add a soft edge effect to a shape."""
    spPr = shape._element.spPr

    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(effectLst)

    softEdge = effectLst.makeelement(qn("a:softEdge"), {
        "rad": str(int(radius_pt * 12700)),
    })
    effectLst.append(softEdge)


def slide1_drop_shadows(prs):
    """Slide 1: Drop shadows with different offsets, blur, and colors."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 1: Drop Shadows (Offset, Blur, Color Variations)", 18)

    configs = [
        ("Small offset\nno blur", 3, 0, 315, (0x00, 0x00, 0x00), 40),
        ("Medium offset\nmedium blur", 6, 4, 315, (0x00, 0x00, 0x00), 50),
        ("Large offset\nlarge blur", 12, 8, 315, (0x00, 0x00, 0x00), 60),
        ("Red shadow", 6, 4, 315, (0xFF, 0x00, 0x00), 60),
        ("Blue shadow", 6, 4, 315, (0x00, 0x00, 0xFF), 60),
    ]

    for i, (label, dist, blur, angle, color, alpha) in enumerate(configs):
        x = 0.4 + (i * 2.5)
        y = 1.0

        shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(2.0), Inches(2.0))
        set_solid_fill(shape, (0x44, 0x72, 0xC4))
        add_drop_shadow(shape, dist, blur, angle, color, alpha)
        add_label(slide, x, y + 2.2, 2.0, 0.6, label, 9)

    # Second row: different directions
    directions = [
        ("Shadow Right (0\u00b0)", 0),
        ("Shadow Bottom-Right (315\u00b0)", 315),
        ("Shadow Bottom (270\u00b0)", 270),
        ("Shadow Left (180\u00b0)", 180),
        ("Shadow Top (90\u00b0)", 90),
    ]

    for i, (label, angle) in enumerate(directions):
        x = 0.4 + (i * 2.5)
        y = 4.0

        shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(2.0), Inches(2.0))
        set_solid_fill(shape, (0xED, 0x7D, 0x31))
        add_drop_shadow(shape, 6, 4, angle, (0x00, 0x00, 0x00), 50)
        add_label(slide, x, y + 2.2, 2.0, 0.5, label, 9)


def slide2_outer_glow(prs):
    """Slide 2: Outer glow effects."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background for glow visibility
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 2: Outer Glow Effects", 18)
    # Re-color the label for dark background
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    glow_configs = [
        ("Small gold glow", 4, (0xFF, 0xD7, 0x00), 50, (0xFF, 0xD7, 0x00)),
        ("Medium red glow", 8, (0xFF, 0x00, 0x00), 50, (0xFF, 0x44, 0x44)),
        ("Large blue glow", 16, (0x00, 0xBF, 0xFF), 60, (0x44, 0x72, 0xC4)),
        ("Green glow", 10, (0x00, 0xFF, 0x00), 50, (0x00, 0x80, 0x00)),
        ("Purple glow", 12, (0xFF, 0x00, 0xFF), 40, (0x80, 0x00, 0x80)),
    ]

    for i, (label, radius, glow_color, alpha, fill_color) in enumerate(glow_configs):
        x = 0.4 + (i * 2.5)
        y = 1.2

        shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(2.0), Inches(2.0))
        set_solid_fill(shape, fill_color)
        add_outer_glow(shape, radius, glow_color, alpha)

        lbl = add_label(slide, x, y + 2.3, 2.0, 0.4, label, 10)
        for p in lbl.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # Second row: glow on different shapes
    shape_configs = [
        ("Rectangle\nwith glow", 1, 10, (0xFF, 0xA5, 0x00), (0xCC, 0x66, 0x00)),
        ("Ellipse\nwith glow", 9, 10, (0x00, 0xFF, 0xFF), (0x00, 0x99, 0x99)),
        ("Rounded rect\nwith glow", 5, 10, (0xFF, 0x69, 0xB4), (0xCC, 0x33, 0x66)),
    ]

    for i, (label, shape_type, radius, glow_color, fill_color) in enumerate(shape_configs):
        x = 1.5 + (i * 4.0)
        y = 4.2

        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(3.0), Inches(2.0))
        set_solid_fill(shape, fill_color)
        add_outer_glow(shape, radius, glow_color, 50)

        lbl = add_label(slide, x, y + 2.2, 3.0, 0.5, label, 10)
        for p in lbl.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)


def slide3_reflection(prs):
    """Slide 3: Reflection effects."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 3: Reflection Effects", 18)

    reflection_configs = [
        ("Tight reflection\n(close, strong)", 0, 70, 30, 0),
        ("Half reflection\n(50% fade)", 0, 50, 50, 0),
        ("Distant reflection\n(offset 10pt)", 0, 40, 40, 10),
        ("Blurred reflection", 3, 50, 50, 2),
        ("Subtle reflection\n(20% alpha)", 0, 20, 30, 0),
    ]

    for i, (label, blur, alpha, end_pos, dist) in enumerate(reflection_configs):
        x = 0.4 + (i * 2.5)
        y = 0.8

        shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(2.0), Inches(2.0))
        set_solid_fill(shape, (0x44, 0x72, 0xC4))
        add_reflection(shape, blur, alpha, end_pos, dist)
        add_label(slide, x, y + 2.8, 2.0, 0.6, label, 9)

    # Second row: reflection on different shapes
    shape_configs = [
        ("Ellipse + reflection", 9, (0xED, 0x7D, 0x31)),
        ("Round rect + reflection", 5, (0x70, 0xAD, 0x47)),
        ("Triangle + reflection", 7, (0xFF, 0xC0, 0x00)),
    ]

    for i, (label, shape_type, color) in enumerate(shape_configs):
        x = 1.0 + (i * 4.2)
        y = 4.2

        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(3.0), Inches(1.8))
        set_solid_fill(shape, color)
        add_reflection(shape, 1, 50, 50, 0)
        add_label(slide, x, y + 2.5, 3.0, 0.4, label, 10)


def slide4_soft_edges(prs):
    """Slide 4: Soft edge effects."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 4: Soft Edge Effects", 18)

    radii = [
        ("2pt soft edge", 2),
        ("5pt soft edge", 5),
        ("10pt soft edge", 10),
        ("20pt soft edge", 20),
        ("40pt soft edge", 40),
    ]

    for i, (label, radius) in enumerate(radii):
        x = 0.4 + (i * 2.5)
        y = 1.0

        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.0), Inches(2.0))
        set_solid_fill(shape, (0xC0, 0x00, 0x00))
        add_soft_edge(shape, radius)
        add_label(slide, x, y + 2.2, 2.0, 0.4, label, 10)

    # Second row: soft edges on different shapes with different fills
    shape_configs = [
        ("Ellipse, 10pt", 9, (0x00, 0x80, 0x00), 10),
        ("Round rect, 15pt", 5, (0x00, 0x00, 0xCC), 15),
        ("Diamond, 10pt", 4, (0xFF, 0xA5, 0x00), 10),
        ("Triangle, 20pt", 7, (0x80, 0x00, 0x80), 20),
    ]

    for i, (label, shape_type, color, radius) in enumerate(shape_configs):
        x = 0.8 + (i * 3.1)
        y = 4.0

        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(2.5), Inches(2.5))
        set_solid_fill(shape, color)
        add_soft_edge(shape, radius)
        add_label(slide, x, y + 2.6, 2.5, 0.4, label, 10)


def slide5_combined_effects(prs):
    """Slide 5: Combined effects on shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 0.3, 0.15, 12, 0.5, "Slide 5: Combined Effects", 18)

    # Shape 1: Shadow + reflection
    shape1 = slide.shapes.add_shape(5, Inches(0.5), Inches(1.0), Inches(3.5), Inches(2.0))
    set_solid_fill(shape1, (0x44, 0x72, 0xC4))
    add_drop_shadow(shape1, 6, 4, 315, (0x00, 0x00, 0x00), 50)
    add_reflection(shape1, 1, 40, 40, 0)
    add_label(slide, 0.5, 3.5, 3.5, 0.4, "Shadow + Reflection", 11)

    # Shape 2: Shadow + soft edge
    shape2 = slide.shapes.add_shape(9, Inches(4.8), Inches(1.0), Inches(3.5), Inches(2.0))
    set_solid_fill(shape2, (0xED, 0x7D, 0x31))
    add_drop_shadow(shape2, 8, 6, 315, (0x00, 0x00, 0x00), 50)
    add_soft_edge(shape2, 8)
    add_label(slide, 4.8, 3.5, 3.5, 0.4, "Shadow + Soft Edge", 11)

    # Shape 3: Glow + reflection
    shape3 = slide.shapes.add_shape(5, Inches(9.0), Inches(1.0), Inches(3.5), Inches(2.0))
    set_solid_fill(shape3, (0x70, 0xAD, 0x47))
    add_outer_glow(shape3, 10, (0x00, 0xFF, 0x00), 40)
    add_reflection(shape3, 1, 50, 50, 0)
    add_label(slide, 9.0, 3.5, 3.5, 0.4, "Glow + Reflection", 11)

    # Shape 4: Shadow + glow + soft edge
    shape4 = slide.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(5.5), Inches(2.0))
    set_solid_fill(shape4, (0x80, 0x00, 0x80))
    add_drop_shadow(shape4, 6, 4, 315, (0x00, 0x00, 0x00), 40)
    add_outer_glow(shape4, 8, (0xFF, 0x00, 0xFF), 30)
    add_soft_edge(shape4, 5)
    add_label(slide, 0.5, 6.7, 5.5, 0.4, "Shadow + Glow + Soft Edge", 11)

    # Shape 5: All effects
    shape5 = slide.shapes.add_shape(9, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.0))
    set_solid_fill(shape5, (0xFF, 0xC0, 0x00))
    add_drop_shadow(shape5, 6, 4, 315, (0x00, 0x00, 0x00), 40)
    add_outer_glow(shape5, 6, (0xFF, 0xA5, 0x00), 30)
    add_reflection(shape5, 0, 30, 30, 0)
    add_soft_edge(shape5, 3)
    add_label(slide, 6.8, 6.7, 5.5, 0.4, "Shadow + Glow + Reflection + Soft Edge", 11)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide1_drop_shadows(prs)
    slide2_outer_glow(prs)
    slide3_reflection(prs)
    slide4_soft_edges(prs)
    slide5_combined_effects(prs)

    output_path = ROOT / "test-data" / "effect-stress-test.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"Generated: {output_path}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
